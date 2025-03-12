import os
import fitz
import docx
import openpyxl
from bs4 import BeautifulSoup
from pptx import Presentation
from typing import List, Tuple
from pathlib import Path
import shutil
import hashlib

from langchain_text_splitters import MarkdownHeaderTextSplitter

headers = [
    ("#", "Header 1"),
    ("##", "Header 2"),
    ("###", "Header 3"),
    ("#####", "Header 4"),
    ("######", "Header 5"),
    ("#######", "Header 6"),
    ("#######", "Header 7"),
]
splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers)

class DocConvertor:
    def __init__(self, image_folder: str = 'images'):
        """初始化文档转换器
        
        Args:
            image_folder: 图片保存的文件夹路径
        """
        self.image_folder = image_folder
        os.makedirs(image_folder, exist_ok=True)
    
    def _save_image(self, image_data: bytes, extension: str = '.png') -> str:
        """保存图片并返回文件路径
        
        Args:
            image_data: 图片二进制数据
            extension: 图片扩展名
            
        Returns:
            str: 保存后的图片路径
        """
        # 使用图片内容的哈希值作为文件名，避免重复
        image_hash = hashlib.md5(image_data).hexdigest()
        image_path = os.path.join(self.image_folder, f'{image_hash}{extension}')
        
        with open(image_path, 'wb') as f:
            f.write(image_data)
        
        return image_path
    
    def convert_docx(self, file_path: str) -> str:
        """将Word文档转换为Markdown
        
        Args:
            file_path: Word文档路径
            
        Returns:
            str: 转换后的Markdown文本
        """
        doc = docx.Document(file_path)
        markdown = []
        
        for para in doc.paragraphs:
            if para.style.name.startswith('Heading'):
                level = para.style.name[-1]
                markdown.append(f"{'#' * int(level)} {para.text}\n")
            else:
                markdown.append(f"{para.text}\n")
        
        # 处理图片
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                image_path = self._save_image(image_data)
                markdown.append(f"![image]({image_path})\n")
        
        return '\n'.join(markdown)
    
    def convert_xlsx(self, file_path: str) -> str:
        """将Excel文件转换为Markdown表格
        
        Args:
            file_path: Excel文件路径
            
        Returns:
            str: 转换后的Markdown文本
        """
        wb = openpyxl.load_workbook(file_path)
        markdown = []
        
        for sheet in wb.worksheets:
            markdown.append(f"## {sheet.title}\n")
            
            # 获取有效区域
            max_row = sheet.max_row
            max_col = sheet.max_column
            
            # 表头
            header = ['|']
            separator = ['|']
            for col in range(1, max_col + 1):
                cell_value = sheet.cell(row=1, column=col).value or ''
                header.append(f" {cell_value} |")
                separator.append(" --- |")
            
            markdown.append(''.join(header))
            markdown.append(''.join(separator))
            
            # 表格内容
            for row in range(2, max_row + 1):
                row_data = ['|']
                for col in range(1, max_col + 1):
                    cell_value = sheet.cell(row=row, column=col).value or ''
                    row_data.append(f" {cell_value} |")
                markdown.append(''.join(row_data))
            markdown.append('\n')
        
        return '\n'.join(markdown)
    
    def convert_pptx(self, file_path: str) -> str:
        """将PPT文件转换为Markdown
        
        Args:
            file_path: PPT文件路径
            
        Returns:
            str: 转换后的Markdown文本
        """
        prs = Presentation(file_path)
        markdown = []
        
        for slide in prs.slides:
            # 处理文本
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    markdown.append(shape.text + '\n')
                
                # 处理图片
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image = shape.image
                    image_path = self._save_image(image.blob)
                    markdown.append(f"![image]({image_path})\n")
        
        return '\n'.join(markdown)
    
    def convert_pdf(self, file_path: str) -> str:
        """将PDF文件转换为Markdown
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            str: 转换后的Markdown文本
        """
        doc = fitz.open(file_path)
        markdown = []
        
        for page in doc:
            # 提取文本
            text = page.get_text()
            markdown.append(text + '\n')
            
            # 提取图片
            for image in page.get_images():
                xref = image[0]
                base_image = doc.extract_image(xref)
                image_data = base_image["image"]
                image_path = self._save_image(image_data)
                markdown.append(f"![image]({image_path})\n")
        
        return '\n'.join(markdown)
    
    def convert_html(self, file_path: str) -> str:
        """将HTML文件转换为Markdown
        
        Args:
            file_path: HTML文件路径
            
        Returns:
            str: 转换后的Markdown文本
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        
        markdown = []
        
        # 处理标题
        for i in range(1, 7):
            for heading in soup.find_all(f'h{i}'):
                markdown.append(f"{'#' * i} {heading.get_text()}\n")
        
        # 处理段落
        for p in soup.find_all('p'):
            markdown.append(f"{p.get_text()}\n")
        
        # 处理图片
        for img in soup.find_all('img'):
            src = img.get('src')
            if src and (src.startswith('http') or src.startswith('https')):
                # 对于网络图片，直接使用URL
                markdown.append(f"![image]({src})\n")
            elif src:
                # 对于本地图片，复制到指定目录
                img_path = os.path.join(os.path.dirname(file_path), src)
                if os.path.exists(img_path):
                    with open(img_path, 'rb') as f:
                        image_data = f.read()
                        new_path = self._save_image(image_data)
                        markdown.append(f"![image]({new_path})\n")
        
        return '\n'.join(markdown)
    
    def convert(self, file_path: str) -> str:
        """根据文件类型调用相应的转换方法
        
        Args:
            file_path: 文件路径
            
        Returns:
            str: 转换后的Markdown文本
        """
        ext = Path(file_path).suffix.lower()
        
        converters = {
            '.docx': self.convert_docx,
            '.xlsx': self.convert_xlsx,
            '.pptx': self.convert_pptx,
            '.pdf': self.convert_pdf,
            '.html': self.convert_html,
            '.htm': self.convert_html
        }
        
        converter = converters.get(ext)
        if not converter:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return converter(file_path)

c = DocConvertor()
text = c.convert("test3.xlsx")
docs = splitter.split_text(text)
print(docs)